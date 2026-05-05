"""Generate a rich, dark-themed PowerPoint for the Two-Tower Rec Sys presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy

# ── Palette ──────────────────────────────────────────────────
BG       = RGBColor(0x1e, 0x1e, 0x2e)   # dark background
SURFACE  = RGBColor(0x31, 0x32, 0x44)   # card surface
PURPLE   = RGBColor(0xcb, 0xa6, 0xf7)   # accent purple
BLUE     = RGBColor(0x89, 0xb4, 0xfa)   # accent blue
GREEN    = RGBColor(0xa6, 0xe3, 0xa1)   # accent green
YELLOW   = RGBColor(0xf9, 0xe2, 0xaf)   # accent yellow
RED      = RGBColor(0xf3, 0x8b, 0xa8)   # accent red
PEACH    = RGBColor(0xfa, 0xb3, 0x87)   # accent peach
WHITE    = RGBColor(0xcd, 0xd6, 0xf4)   # near-white text
SUBTEXT  = RGBColor(0xa6, 0xad, 0xc8)   # muted text
OVERLAY  = RGBColor(0x45, 0x47, 0x5a)   # overlay surface

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=SURFACE, line=None, line_w=Pt(1.5), alpha=None):
    """Add a filled rectangle."""
    shp = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def label_box(slide, label, l, t, w, h, fill=SURFACE, border=PURPLE,
              font_size=14, bold=True, color=WHITE, sub=None, sub_color=None):
    """Rounded-looking card with optional subtitle."""
    b = box(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(2))
    # title
    txb = slide.shapes.add_textbox(l, t, w, h * (0.55 if sub else 1.0))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if sub:
        stxb = slide.shapes.add_textbox(l, t + h * 0.52, w, h * 0.48)
        stf = stxb.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = sub
        srun.font.size = Pt(font_size - 3)
        srun.font.color.rgb = sub_color or SUBTEXT
    return b

def arrow_down(slide, cx, y, h=Inches(0.3), color=SUBTEXT):
    """Vertical down arrow."""
    shp = slide.shapes.add_connector(1, cx, y, cx, y + h)
    shp.line.color.rgb = color
    shp.line.width = Pt(2)

def arrow_right(slide, x, cy, w=Inches(0.3), color=SUBTEXT):
    shp = slide.shapes.add_connector(1, x, cy, x + w, cy)
    shp.line.color.rgb = color
    shp.line.width = Pt(2)

def chip(slide, text, l, t, w=Inches(1.4), h=Inches(0.35),
         fill=PURPLE, color=BG, size=11, bold=True):
    """Small pill chip."""
    b = box(slide, l, t, w, h, fill=fill)
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def section_header(slide, text, sub=None):
    """Top section header bar."""
    box(slide, 0, 0, W, Inches(1.1), fill=SURFACE)
    txt(slide, text, Inches(0.5), Inches(0.12), Inches(10), Inches(0.55),
        size=28, bold=True, color=PURPLE)
    if sub:
        txt(slide, sub, Inches(0.5), Inches(0.6), Inches(12), Inches(0.4),
            size=14, color=SUBTEXT)
    # accent bar
    box(slide, 0, Inches(1.1), W, Inches(0.04), fill=PURPLE)

def stat_card(slide, value, label, l, t, w=Inches(1.7), h=Inches(1.0),
              val_color=PURPLE):
    box(slide, l, t, w, h, fill=OVERLAY, line=SURFACE)
    txt(slide, value, l, t + Inches(0.05), w, Inches(0.55),
        size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + Inches(0.55), w, Inches(0.4),
        size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

def add_hyperlink(shape, url):
    """Add a click-through hyperlink to a shape."""
    click = shape.click_action
    click.hyperlink.address = url


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — HOOK / COVER
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)

# Big dramatic title
txt(s, "You open Netflix.", Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
    size=40, bold=True, color=WHITE)
txt(s, "It already knows what you want to watch.", Inches(0.8), Inches(1.35), Inches(11), Inches(0.7),
    size=32, color=SUBTEXT)
txt(s, "That's not magic.", Inches(0.8), Inches(2.1), Inches(11), Inches(0.6),
    size=32, bold=True, color=PURPLE)

# Divider
box(s, Inches(0.8), Inches(2.85), Inches(4), Inches(0.05), fill=PURPLE)

txt(s, "I built one from scratch.", Inches(0.8), Inches(3.05), Inches(11), Inches(0.55),
    size=22, color=WHITE)
txt(s, "4 deep learning models  ·  98,000 real users  ·  12 ablation variants",
    Inches(0.8), Inches(3.55), Inches(11), Inches(0.5), size=16, color=SUBTEXT)
txt(s, "Figured out what actually works — and what doesn't.",
    Inches(0.8), Inches(4.0), Inches(11), Inches(0.5), size=18, color=GREEN, bold=True)

# Bottom right: name + course
txt(s, "Nidhi Rajani  |  EAS 509  |  Spring 2026",
    Inches(8.5), Inches(6.9), Inches(4.5), Inches(0.4), size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)

# Corner chip
chip(s, "Amazon Video Games 2023", Inches(0.8), Inches(6.85), w=Inches(2.8), h=Inches(0.4),
     fill=OVERLAY, color=YELLOW, size=10, bold=False)

# Floating company logos text
txt(s, "Used by:  YouTube  ·  Pinterest  ·  DoorDash  ·  Airbnb  ·  Spotify  ·  Twitter/X",
    Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.45), size=13, color=OVERLAY)
txt(s, "Used by:  YouTube  ·  Pinterest  ·  DoorDash  ·  Airbnb  ·  Spotify  ·  Twitter/X",
    Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.45), size=13, color=SUBTEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "The Problem", "Why recommendation is hard")

# Left: sparsity visual
box(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.7), fill=OVERLAY)
txt(s, "User–Item Matrix", Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.4),
    size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# Draw grid cells — mostly empty, a few purple
import random; random.seed(42)
cell_w = Inches(0.28); cell_h = Inches(0.28)
cols, rows = 14, 14
ox = Inches(0.65); oy = Inches(1.85)
filled = set(random.sample(range(cols*rows), 6))
for r in range(rows):
    for c in range(cols):
        idx = r*cols + c
        fc = PURPLE if idx in filled else RGBColor(0x31,0x32,0x44)
        b = box(s, ox + c*(cell_w + Inches(0.02)),
                   oy + r*(cell_h + Inches(0.02)),
                   cell_w, cell_h, fill=fc, line=SURFACE, line_w=Pt(0.5))

txt(s, "99.97% EMPTY", Inches(0.5), Inches(6.3), Inches(5.5), Inches(0.5),
    size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s, "Only 6 of ~196 cells filled above", Inches(0.5), Inches(6.75), Inches(5.5), Inches(0.35),
    size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Right: three problem cards
problems = [
    ("🕸  Sparsity", "Most users bought only 5-7 games.\nMost items have very few reviews.\nHardly any signal to learn from.", RED),
    ("👤  Cold-Start", "A brand-new user arrives.\nNo purchase history.\nMF and graph models output nothing.", YELLOW),
    ("⚡  Scale", "98,906 users × 26,354 items\n= 2.6 billion dot products per request.\nNeeds to respond in microseconds.", BLUE),
]
for i, (title, body, col) in enumerate(problems):
    ty = Inches(1.4) + i * Inches(2.0)
    box(s, Inches(6.3), ty, Inches(6.5), Inches(1.7), fill=SURFACE, line=col, line_w=Pt(2))
    txt(s, title, Inches(6.5), ty + Inches(0.12), Inches(6), Inches(0.4),
        size=15, bold=True, color=col)
    txt(s, body, Inches(6.5), ty + Inches(0.5), Inches(6.0), Inches(1.1),
        size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Dataset — Amazon Video Games 2023", "McAuley Lab, UCSD  ·  amazon-reviews-2023.github.io")

stats = [
    ("98,906", "Users"),
    ("26,354", "Items"),
    ("659K",   "Interactions"),
    ("99.97%", "Sparsity"),
    ("8",      "User Features"),
    ("15",     "Item Features"),
    ("384d",   "Text Embeddings"),
]
sw = Inches(1.65)
for i, (val, lbl) in enumerate(stats):
    stat_card(s, val, lbl, Inches(0.4) + i*(sw + Inches(0.1)), Inches(1.3),
              w=sw, val_color=PURPLE if i < 4 else GREEN)

# Pipeline
box(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.05), fill=PURPLE)
txt(s, "Data Pipeline", Inches(0.5), Inches(2.7), Inches(6), Inches(0.4),
    size=14, bold=True, color=PURPLE)

steps = [
    ("1. Download", "Amazon 2023\nHuggingFace"),
    ("2. K-core Filter", "k=5\nMin interactions"),
    ("3. Feature Eng.", "8 user feats\n15 item feats"),
    ("4. Text Embed", "SentenceTrans\nall-MiniLM-L6"),
    ("5. Train Split", "Leave-last-2\nchronological"),
    ("6. Train Models", "MF · LightGCN\nTT · FG-LightGCN"),
]
sw2 = Inches(1.95)
for i, (step, detail) in enumerate(steps):
    lx = Inches(0.4) + i*(sw2 + Inches(0.1))
    label_box(s, step, lx, Inches(3.2), sw2, Inches(1.1),
              fill=OVERLAY, border=BLUE, font_size=12, sub=detail, sub_color=SUBTEXT)
    if i < len(steps)-1:
        arrow_right(s, lx + sw2, Inches(3.2) + Inches(0.55), w=Inches(0.1))

# Features detail
box(s, Inches(0.5), Inches(4.55), Inches(5.8), Inches(2.5), fill=OVERLAY, line=BLUE, line_w=Pt(1.5))
txt(s, "User Features (8)", Inches(0.6), Inches(4.65), Inches(5.5), Inches(0.4),
    size=13, bold=True, color=BLUE)
user_feats = "• Total interactions  • Avg rating given  • Rating variance\n• Tenure (days active)  • Recency  • Avg price paid\n• 30-day activity count  • Category breadth"
txt(s, user_feats, Inches(0.6), Inches(5.1), Inches(5.5), Inches(1.8), size=11, color=WHITE)

box(s, Inches(6.6), Inches(4.55), Inches(6.2), Inches(2.5), fill=OVERLAY, line=GREEN, line_w=Pt(1.5))
txt(s, "Item Features (15) + Text", Inches(6.7), Inches(4.65), Inches(5.9), Inches(0.4),
    size=13, bold=True, color=GREEN)
item_feats = "• Price  • Avg rating  • Rating count  • Rating variance\n• Category (encoded)  • Has description flag\n• SentenceTransformer title embeddings → 384d\n  all-MiniLM-L6-v2  (22M param distilled BERT)"
txt(s, item_feats, Inches(6.7), Inches(5.1), Inches(5.9), Inches(1.8), size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — RESEARCH FOUNDATION
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Research Foundation", "Papers that this system builds on")

papers = [
    {
        "title": "Matrix Factorization Techniques\nfor Recommender Systems",
        "authors": "Koren, Bell, Volinsky  —  IEEE Computer 2009",
        "venue": "IEEE Computer",
        "what": "Introduced the embedding factorization idea.\nEvery user and item gets a latent vector.\nDot product = predicted rating.",
        "used": "Direct baseline. BPR loss variant (Rendle 2009).",
        "color": SUBTEXT,
    },
    {
        "title": "Neural Collaborative Filtering",
        "authors": "He et al.  —  WWW 2017",
        "venue": "WWW 2017",
        "what": "Replace dot product with MLP.\nShowed neural nets can model non-linear user-item interactions.",
        "used": "Motivation for MLP layers in the Two-Tower towers.",
        "color": BLUE,
    },
    {
        "title": "LightGCN: Simplifying and\nPowering Graph Convolution",
        "authors": "He et al.  —  SIGIR 2020",
        "venue": "SIGIR 2020",
        "what": "Remove feature transforms + nonlinearities from GCN.\nPure neighborhood averaging on interaction graph.\nMean pool over layers for multi-scale representations.",
        "used": "Directly implemented. Best accuracy model (HR@10=0.729).",
        "color": PURPLE,
    },
    {
        "title": "Deep Neural Networks for\nYouTube Recommendations",
        "authors": "Covington, Adams, Sargin  —  RecSys 2016",
        "venue": "RecSys 2016",
        "what": "Two-stage: retrieval (Two-Tower + ANN)\nthen ranking (deep network).\nIn-batch negative sampling trick.",
        "used": "Architecture foundation for Two-Tower v5.\nImproved with GRU, InfoNCE, FAISS HNSW.",
        "color": GREEN,
    },
]

pw = Inches(3.05)
for i, p in enumerate(papers):
    lx = Inches(0.3) + i*(pw + Inches(0.1))
    box(s, lx, Inches(1.3), pw, Inches(5.7), fill=SURFACE, line=p["color"], line_w=Pt(2))
    chip(s, p["venue"], lx + Inches(0.1), Inches(1.38), w=pw - Inches(0.2), h=Inches(0.3),
         fill=p["color"], color=BG, size=9)
    txt(s, p["title"], lx + Inches(0.12), Inches(1.78), pw - Inches(0.2), Inches(0.85),
        size=12, bold=True, color=p["color"])
    txt(s, p["authors"], lx + Inches(0.12), Inches(2.65), pw - Inches(0.2), Inches(0.4),
        size=9, color=SUBTEXT, italic=True)
    box(s, lx + Inches(0.1), Inches(3.1), pw - Inches(0.2), Inches(0.03), fill=p["color"])
    txt(s, "What it showed:", lx + Inches(0.12), Inches(3.18), pw - Inches(0.2), Inches(0.3),
        size=10, bold=True, color=WHITE)
    txt(s, p["what"], lx + Inches(0.12), Inches(3.48), pw - Inches(0.2), Inches(1.4),
        size=10, color=WHITE)
    txt(s, "How I used it:", lx + Inches(0.12), Inches(4.95), pw - Inches(0.2), Inches(0.3),
        size=10, bold=True, color=p["color"])
    txt(s, p["used"], lx + Inches(0.12), Inches(5.28), pw - Inches(0.2), Inches(0.85),
        size=10, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — MF ARCHITECTURE
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Model 1: Matrix Factorization", "Baseline  ·  BPR Loss  ·  HR@10 = 0.6825")

chip(s, "BASELINE", Inches(11.5), Inches(0.15), w=Inches(1.5), h=Inches(0.35),
     fill=SUBTEXT, color=BG, size=10)

# Architecture diagram
cx = Inches(4.5)

label_box(s, "User ID", Inches(1.0), Inches(1.7), Inches(2.2), Inches(0.8),
          fill=SURFACE, border=PURPLE, font_size=14, sub="e.g. User #100", sub_color=SUBTEXT)
label_box(s, "Item ID", Inches(9.8), Inches(1.7), Inches(2.2), Inches(0.8),
          fill=SURFACE, border=BLUE, font_size=14, sub="e.g. Game #452", sub_color=SUBTEXT)

arrow_down(s, Inches(2.1), Inches(2.5), h=Inches(0.5))
arrow_down(s, Inches(10.9), Inches(2.5), h=Inches(0.5))

label_box(s, "User Embedding", Inches(0.8), Inches(3.0), Inches(2.6), Inches(1.0),
          fill=OVERLAY, border=PURPLE, font_size=14, bold=True, color=PURPLE,
          sub="64 dimensions", sub_color=WHITE)
label_box(s, "Item Embedding", Inches(9.6), Inches(3.0), Inches(2.6), Inches(1.0),
          fill=OVERLAY, border=BLUE, font_size=14, bold=True, color=BLUE,
          sub="64 dimensions", sub_color=WHITE)

# converging arrows
arrow_right(s, Inches(3.4), Inches(3.5), w=Inches(2.1))
arrow_right(s, Inches(7.5), Inches(3.5), w=Inches(2.1))

label_box(s, "Dot Product → Score", Inches(5.5), Inches(2.9), Inches(2.5), Inches(1.2),
          fill=OVERLAY, border=GREEN, font_size=14, bold=True, color=GREEN,
          sub="sim(u,i) = eᵤ · eᵢ", sub_color=WHITE)

arrow_down(s, Inches(6.75), Inches(4.1), h=Inches(0.5))

label_box(s, "BPR Loss", Inches(5.3), Inches(4.6), Inches(3.0), Inches(1.4),
          fill=OVERLAY, border=PEACH, font_size=16, bold=True, color=PEACH,
          sub="−log σ(s_pos − s_neg)", sub_color=WHITE)

# Insight box
box(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.0), fill=OVERLAY, line=SUBTEXT, line_w=Pt(1))
txt(s, "💡  No features. No graph. No text.  Just ID co-occurrence.  "
       "The floor — anything fancier must beat 0.6825 to justify complexity.",
    Inches(0.7), Inches(6.2), Inches(12), Inches(0.7), size=13, color=WHITE)

# Right side: formula
txt(s, "Parameters:\n~8 Million\n(98K + 26K) × 64",
    Inches(9.5), Inches(5.8), Inches(3.5), Inches(1.5), size=12, color=SUBTEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — LIGHTGCN ARCHITECTURE
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Model 2: LightGCN", "Graph Neural Network  ·  Best Accuracy  ·  HR@10 = 0.7290")

chip(s, "ACCURACY WINNER", Inches(10.8), Inches(0.15), w=Inches(2.2), h=Inches(0.35),
     fill=GREEN, color=BG, size=10)

# Left: graph visual
box(s, Inches(0.4), Inches(1.3), Inches(4.8), Inches(5.7), fill=OVERLAY)
txt(s, "Bipartite Interaction Graph", Inches(0.4), Inches(1.35), Inches(4.8), Inches(0.4),
    size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# User nodes (left col)
for i, (y, label) in enumerate([(Inches(2.0), "User A"), (Inches(3.3), "User B"), (Inches(4.6), "User C")]):
    box(s, Inches(0.7), y, Inches(1.2), Inches(0.55), fill=SURFACE, line=PURPLE, line_w=Pt(2))
    txt(s, label, Inches(0.7), y, Inches(1.2), Inches(0.55), size=11, bold=True,
        color=PURPLE, align=PP_ALIGN.CENTER)

# Item nodes (right col)
for i, (y, label) in enumerate([(Inches(1.6), "Game 1"), (Inches(2.7), "Game 2"),
                                  (Inches(3.8), "Game 3"), (Inches(4.9), "Game 4")]):
    box(s, Inches(3.7), y, Inches(1.2), Inches(0.55), fill=SURFACE, line=BLUE, line_w=Pt(2))
    txt(s, label, Inches(3.7), y, Inches(1.2), Inches(0.55), size=11, bold=True,
        color=BLUE, align=PP_ALIGN.CENTER)

# Edges (connections)
edges = [(Inches(1.9), Inches(2.27), Inches(3.7), Inches(1.87)),   # A→1
         (Inches(1.9), Inches(2.27), Inches(3.7), Inches(2.97)),   # A→2
         (Inches(1.9), Inches(3.57), Inches(3.7), Inches(2.97)),   # B→2
         (Inches(1.9), Inches(3.57), Inches(3.7), Inches(4.07)),   # B→3
         (Inches(1.9), Inches(4.87), Inches(3.7), Inches(4.07)),   # C→3
         (Inches(1.9), Inches(4.87), Inches(3.7), Inches(5.17)),]  # C→4
for x1,y1,x2,y2 in edges:
    c = s.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = SUBTEXT
    c.line.width = Pt(1.5)

txt(s, "Edges = purchases\n(659K total)", Inches(0.5), Inches(6.0), Inches(4.6), Inches(0.6),
    size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Right: propagation layers
lw = Inches(5.8)
lx = Inches(5.7)
layers = [
    ("Layer 0", "Raw ID Embeddings", "98K users + 26K items → 64d each", BLUE),
    ("Layer 1", "Direct Neighbors", "User A's emb = avg of Game 1, Game 2 embs", BLUE),
    ("Layer 2", "2-Hop Neighbors", "'Friends-of-friends' purchase patterns", BLUE),
    ("Layer 3", "3-Hop Neighborhood", "Deep transitive community signals", BLUE),
    ("MEAN POOL", "Final Embedding", "E_final = mean(E⁰, E¹, E², E³)", GREEN),
]
for i, (tag, title, detail, col) in enumerate(layers):
    ly = Inches(1.4) + i * Inches(1.18)
    box(s, lx, ly, lw, Inches(1.0), fill=SURFACE, line=col, line_w=Pt(2))
    chip(s, tag, lx + Inches(0.1), ly + Inches(0.08), w=Inches(1.3), h=Inches(0.28),
         fill=col, color=BG, size=9)
    txt(s, title, lx + Inches(1.5), ly + Inches(0.08), lw - Inches(1.6), Inches(0.35),
        size=13, bold=True, color=col)
    txt(s, detail, lx + Inches(0.15), ly + Inches(0.55), lw - Inches(0.3), Inches(0.38),
        size=11, color=WHITE)
    if i < len(layers)-1:
        txt(s, "↓  Â·E (sparse mm)", lx + Inches(1.8), ly + Inches(1.03),
            Inches(3), Inches(0.3), size=10, color=SUBTEXT)

# Key formula
box(s, lx, Inches(7.1), lw, Inches(0.25), fill=OVERLAY)

# Fatal flaw
box(s, Inches(0.4), Inches(6.85), Inches(4.8), Inches(0.55), fill=SURFACE, line=RED, line_w=Pt(1.5))
txt(s, "⚠  Cannot handle cold-start or FAISS serving", Inches(0.5), Inches(6.88), Inches(4.6), Inches(0.45),
    size=11, color=RED, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — TWO-TOWER ARCHITECTURE
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Model 3: Two-Tower", "Production Retrieval  ·  YouTube · Pinterest · DoorDash  ·  HR@10 = 0.6395")

chip(s, "PRODUCTION STANDARD", Inches(10.2), Inches(0.15), w=Inches(2.8), h=Inches(0.35),
     fill=PURPLE, color=BG, size=10)

# USER TOWER
box(s, Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.3), fill=SURFACE, line=PURPLE, line_w=Pt(2))
txt(s, "USER TOWER", Inches(0.3), Inches(1.35), Inches(4.5), Inches(0.45),
    size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

u_inputs = [
    ("User ID Embedding", "64d  ·  Learnable per user", PURPLE),
    ("GRU Sequence Encoder", "Last 20 items → 64d  ·  Temporal patterns", PURPLE),
    ("User Features → Linear", "8 feats → 64d  ·  Activity, recency…", PURPLE),
]
for i, (title, detail, col) in enumerate(u_inputs):
    iy = Inches(1.9) + i * Inches(1.05)
    box(s, Inches(0.45), iy, Inches(4.2), Inches(0.85), fill=OVERLAY, line=col, line_w=Pt(1.5))
    txt(s, title, Inches(0.55), iy + Inches(0.05), Inches(4.0), Inches(0.38),
        size=12, bold=True, color=col)
    txt(s, detail, Inches(0.55), iy + Inches(0.44), Inches(4.0), Inches(0.36),
        size=10, color=WHITE)
    if i < 2:
        txt(s, "+", Inches(2.1), iy + Inches(0.88), Inches(0.4), Inches(0.3),
            size=16, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

txt(s, "↓  Concat → MLP → LayerNorm", Inches(0.5), Inches(5.1), Inches(4.2), Inches(0.35),
    size=10, color=SUBTEXT)
box(s, Inches(0.45), Inches(5.45), Inches(4.2), Inches(0.75), fill=PURPLE, line=PURPLE)
txt(s, "L2-Normalized 64d Vector", Inches(0.45), Inches(5.45), Inches(4.2), Inches(0.75),
    size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)

# ITEM TOWER
box(s, Inches(8.5), Inches(1.3), Inches(4.5), Inches(5.3), fill=SURFACE, line=BLUE, line_w=Pt(2))
txt(s, "ITEM TOWER", Inches(8.5), Inches(1.35), Inches(4.5), Inches(0.45),
    size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

i_inputs = [
    ("Item ID Embedding", "64d  ·  Learnable per item", BLUE),
    ("Text Embedding (SentenceTrans)", "Title → 384d → Linear → 64d", BLUE),
    ("Item Features → Linear", "15 feats → 64d  ·  Price, rating…", BLUE),
]
for i, (title, detail, col) in enumerate(i_inputs):
    iy = Inches(1.9) + i * Inches(1.05)
    box(s, Inches(8.65), iy, Inches(4.2), Inches(0.85), fill=OVERLAY, line=col, line_w=Pt(1.5))
    txt(s, title, Inches(8.75), iy + Inches(0.05), Inches(4.0), Inches(0.38),
        size=12, bold=True, color=col)
    txt(s, detail, Inches(8.75), iy + Inches(0.44), Inches(4.0), Inches(0.36),
        size=10, color=WHITE)
    if i < 2:
        txt(s, "+", Inches(10.4), iy + Inches(0.88), Inches(0.4), Inches(0.3),
            size=16, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

txt(s, "↓  Concat → MLP → LayerNorm", Inches(8.65), Inches(5.1), Inches(4.2), Inches(0.35),
    size=10, color=SUBTEXT)
box(s, Inches(8.65), Inches(5.45), Inches(4.2), Inches(0.75), fill=BLUE, line=BLUE)
txt(s, "L2-Normalized 64d Vector", Inches(8.65), Inches(5.45), Inches(4.2), Inches(0.75),
    size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)

# Center: dot product
txt(s, "→", Inches(4.85), Inches(5.68), Inches(0.4), Inches(0.4),
    size=22, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)
box(s, Inches(5.2), Inches(5.3), Inches(2.9), Inches(1.15), fill=OVERLAY, line=GREEN, line_w=Pt(2))
txt(s, "Dot Product", Inches(5.2), Inches(5.35), Inches(2.9), Inches(0.4),
    size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "InfoNCE Loss  τ=0.2", Inches(5.2), Inches(5.75), Inches(2.9), Inches(0.35),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "←", Inches(8.1), Inches(5.68), Inches(0.4), Inches(0.4),
    size=22, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Cold start callout
box(s, Inches(4.8), Inches(1.3), Inches(3.6), Inches(3.6), fill=OVERLAY, line=GREEN, line_w=Pt(1.5))
txt(s, "✅ COLD-START", Inches(4.85), Inches(1.38), Inches(3.5), Inches(0.38),
    size=13, bold=True, color=GREEN)
txt(s, "New user? No ID needed.\nGRU runs on browsed\nitem text embeddings.\nWorks from 3 items.", Inches(4.85), Inches(1.78), Inches(3.5), Inches(1.6),
    size=12, color=WHITE)
box(s, Inches(4.85), Inches(3.35), Inches(3.45), Inches(1.35), fill=SURFACE, line=GREEN, line_w=Pt(1))
txt(s, "⚡ FAISS HNSW\n29 μs / query\n34K queries/sec", Inches(4.85), Inches(3.38), Inches(3.45), Inches(1.2),
    size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — TWO-TOWER vs YOUTUBE ORIGINAL
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "What I Improved Over the Original", "YouTube 2016 Two-Tower → My Two-Tower v5")

rows = [
    ("Component", "YouTube 2016 Original", "My Two-Tower v5", "Impact"),
    ("User Signal", "Watch history avg pooling\n(order doesn't matter)", "GRU over last 20 items\n(sequential — order matters)", "+0.4% HR@10\nvs avg pool baseline"),
    ("Text", "Not used", "SentenceTransformer 384d\nall-MiniLM-L6-v2 → 64d", "+2.6% HR@10\nvs no-text baseline"),
    ("Loss Function", "Softmax over all items\n(slow, needs full catalog)", "InfoNCE with in-batch\nnegatives, τ=0.2", "10× faster training\nstronger gradients"),
    ("Serving", "Nearest neighbor lookup\n(no specific index)", "FAISS HNSW index\nOffline pre-computation", "29μs per query\n34K qps"),
    ("Cold-Start", "Not addressed at all", "GRU encodes browsed\nitem text — no ID needed", "Only model that\nworks for new users"),
]

col_widths = [Inches(2.0), Inches(3.2), Inches(3.2), Inches(2.5)]
ty = Inches(1.3)
rh = Inches(0.88)
header_colors = [OVERLAY, SUBTEXT, PURPLE, GREEN]
for ri, row in enumerate(rows):
    for ci, (cell, cw) in enumerate(zip(row, col_widths)):
        lx = Inches(0.4) + sum(col_widths[:ci])
        is_header = ri == 0
        fc = OVERLAY if is_header else (SURFACE if ri % 2 == 0 else OVERLAY)
        lc = header_colors[ci] if is_header else SURFACE
        box(s, lx, ty + ri*rh, cw - Inches(0.05), rh - Inches(0.05),
            fill=fc, line=lc if is_header else None, line_w=Pt(1.5))
        tc = [SUBTEXT, WHITE, PURPLE, GREEN][ci] if is_header else WHITE
        txt(s, cell, lx + Inches(0.08), ty + ri*rh + Inches(0.05),
            cw - Inches(0.2), rh - Inches(0.1),
            size=11 if ri > 0 else 12, bold=(ri==0), color=tc)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — FEATURE-GATED LIGHTGCN (MY CONTRIBUTION)
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Model 4: Feature-Gated LightGCN", "My Novel Contribution  ·  HR@10 = 0.7190  ·  Gate = 0.18")

chip(s, "MY CONTRIBUTION", Inches(10.5), Inches(0.15), w=Inches(2.5), h=Inches(0.35),
     fill=YELLOW, color=BG, size=10)

# Research question
box(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.75), fill=OVERLAY, line=YELLOW, line_w=Pt(1.5))
txt(s, "Research question:  Does adding side features (user stats, item stats, text embeddings) to LightGCN improve recommendations?",
    Inches(0.55), Inches(1.38), Inches(12.1), Inches(0.6), size=13, bold=True, color=YELLOW)

# Architecture
bw = Inches(3.8)
bh = Inches(1.4)
by1 = Inches(2.25)
box(s, Inches(0.5), by1, bw, bh, fill=SURFACE, line=BLUE, line_w=Pt(2))
txt(s, "LightGCN (3 layers)", Inches(0.6), by1+Inches(0.1), bw-Inches(0.2), Inches(0.4),
    size=14, bold=True, color=BLUE)
txt(s, "Graph propagation\nÂ·E × 3 layers → mean pool\n→ 64d E_graph", Inches(0.6), by1+Inches(0.5), bw-Inches(0.2), Inches(0.8),
    size=11, color=WHITE)

box(s, Inches(4.8), by1, bw, bh, fill=SURFACE, line=YELLOW, line_w=Pt(2))
txt(s, "Feature Projections", Inches(4.9), by1+Inches(0.1), bw-Inches(0.2), Inches(0.4),
    size=14, bold=True, color=YELLOW)
txt(s, "Linear(user_feats 8→64)\nLinear(item_feats 15→64)\nLinear(text 384→64)", Inches(4.9), by1+Inches(0.5), bw-Inches(0.2), Inches(0.8),
    size=11, color=WHITE)

# Gate
arrow_down(s, Inches(2.4), by1+bh, h=Inches(0.35))
arrow_down(s, Inches(6.7), by1+bh, h=Inches(0.35))
box(s, Inches(2.5), Inches(4.1), Inches(5.0), Inches(1.5), fill=OVERLAY, line=RED, line_w=Pt(2.5))
txt(s, "Learnable Sigmoid Gate", Inches(2.5), Inches(4.18), Inches(5.0), Inches(0.4),
    size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s, "gate = sigmoid(θ)    [single learnable parameter]", Inches(2.55), Inches(4.6), Inches(4.9), Inches(0.35),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "E_final = (1−gate)·E_graph + gate·E_features", Inches(2.55), Inches(4.95), Inches(4.9), Inches(0.35),
    size=11, color=WHITE, align=PP_ALIGN.CENTER)

arrow_down(s, Inches(5.0), Inches(5.6), h=Inches(0.35))
box(s, Inches(3.5), Inches(5.95), Inches(3.0), Inches(0.75), fill=GREEN, line=GREEN)
txt(s, "BPR Loss → Optimize", Inches(3.5), Inches(5.95), Inches(3.0), Inches(0.75),
    size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)

# Gate convergence result
box(s, Inches(9.1), Inches(2.2), Inches(3.9), Inches(5.1), fill=SURFACE, line=RED, line_w=Pt(2))
txt(s, "Gate Convergence", Inches(9.2), Inches(2.3), Inches(3.7), Inches(0.4),
    size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s, "Started at: 0.57", Inches(9.2), Inches(2.78), Inches(3.7), Inches(0.35),
    size=12, color=SUBTEXT, align=PP_ALIGN.CENTER)
txt(s, "Converged to: 0.18", Inches(9.2), Inches(3.15), Inches(3.7), Inches(0.35),
    size=12, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(9.2), Inches(3.55), Inches(3.6), Inches(0.65), fill=BLUE, line=BLUE)
txt(s, "82% graph signal", Inches(9.2), Inches(3.58), Inches(3.6), Inches(0.58),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(9.2), Inches(4.25), Inches(3.6), Inches(0.65), fill=YELLOW, line=YELLOW)
txt(s, "18% features", Inches(9.2), Inches(4.28), Inches(3.6), Inches(0.58),
    size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
txt(s, "Model self-learned:\ngraph is 4.5× more\nvaluable than features\non 99.97% sparse data",
    Inches(9.2), Inches(5.0), Inches(3.7), Inches(1.1), size=12, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "✅ Confirms ablation study", Inches(9.2), Inches(6.1), Inches(3.7), Inches(0.35),
    size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — LIVE DEMO
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Live Demo", "two-tower-rec-sys.streamlit.app")

# Giant demo link card
box(s, Inches(1.5), Inches(1.4), Inches(10.3), Inches(2.0), fill=SURFACE, line=PURPLE, line_w=Pt(3))
txt(s, "🚀  Open Live Demo", Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.65),
    size=22, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
url_shape = txt(s, "https://two-tower-rec-sys.streamlit.app",
                Inches(1.5), Inches(2.1), Inches(10.3), Inches(0.6),
                size=18, color=BLUE, align=PP_ALIGN.CENTER, bold=True)

# Demo what to show
demos = [
    ("Demo 1: Existing User", "MF vs Two-Tower", "• Pick User ID 100\n• See purchase history\n• Compare 10 recs side-by-side\n• Note 4-5 overlapping items\n  (high confidence from both)", PURPLE),
    ("Demo 2: Cold-Start", "New User — Zero History", "• Select 'Souls-like' scenario\n• Dark Souls · Elden Ring · Sekiro\n• GRU encodes text embeddings\n• FAISS retrieves in milliseconds\n• MF: ✗   LightGCN: ✗   TT: ✅", GREEN),
    ("Demo 3: Architecture", "Interactive Explorer", "• Navigate to Model Architectures\n• Hover over each box\n• See details on hover\n• Show FG-LightGCN tab\n• Show gate convergence chart", BLUE),
]
for i, (title, sub, body, col) in enumerate(demos):
    lx = Inches(0.4) + i * Inches(4.3)
    box(s, lx, Inches(3.7), Inches(4.1), Inches(3.4), fill=SURFACE, line=col, line_w=Pt(2))
    chip(s, sub, lx + Inches(0.1), Inches(3.78), w=Inches(3.9), h=Inches(0.3),
         fill=col, color=BG, size=9)
    txt(s, title, lx + Inches(0.15), Inches(4.12), Inches(3.8), Inches(0.4),
        size=14, bold=True, color=col)
    txt(s, body, lx + Inches(0.15), Inches(4.55), Inches(3.8), Inches(2.3),
        size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Results", "Sampled Eval (100 negatives) + Full Ranking (all 26,354 items)")

# Bar chart via pptx
chart_data = ChartData()
chart_data.categories = ["MF (BPR)", "Two-Tower v5", "FG-LightGCN", "LightGCN"]
chart_data.add_series("HR@10 (Sampled)", (0.6825, 0.6395, 0.7190, 0.7290))
chart_data.add_series("NDCG@10 (Sampled)", (0.4550, 0.4148, 0.4824, 0.4940))

chart = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.3), Inches(6.2), Inches(4.5),
    chart_data
).chart

chart.has_title = True
chart.chart_title.text_frame.text = "Sampled Evaluation"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)

plot = chart.plots[0]
series_colors = [PURPLE, GREEN]
for i, series in enumerate(plot.series):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = series_colors[i]

chart.font.color.rgb = WHITE
chart.font.size = Pt(10)

try:
    chart.plot_area.format.fill.background()
    chart.chart_area.format.fill.solid()
    chart.chart_area.format.fill.fore_color.rgb = SURFACE
except: pass

# Full ranking numbers table
box(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(2.8), fill=SURFACE, line=SUBTEXT, line_w=Pt(1))
txt(s, "Full Ranking (Publication Standard — all 26K items)",
    Inches(6.9), Inches(1.38), Inches(5.9), Inches(0.4), size=12, bold=True, color=WHITE)
full_rows = [
    ("Model", "HR@10", "HR@20", "NDCG@10"),
    ("MF (BPR)",     "0.0420", "0.0610", "0.0222"),
    ("Two-Tower v5", "0.0240", "0.0430", "0.0113"),
    ("LightGCN",     "0.0430", "0.0720", "0.0247"),
]
for ri, row in enumerate(full_rows):
    for ci, cell in enumerate(row):
        lx = Inches(6.85) + ci*Inches(1.45)
        ly = Inches(1.82) + ri*Inches(0.5)
        fc = OVERLAY if ri == 0 else (SURFACE if ri%2==0 else OVERLAY)
        box(s, lx, ly, Inches(1.4), Inches(0.45), fill=fc)
        color = WHITE if ri == 0 else (GREEN if (ri==3 and ci>0) else WHITE)
        bold = ri==0 or (ri==3 and ci>0)
        txt(s, cell, lx+Inches(0.05), ly+Inches(0.05), Inches(1.35), Inches(0.38),
            size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)

# Cold start table
box(s, Inches(6.8), Inches(4.3), Inches(6.1), Inches(2.8), fill=SURFACE, line=GREEN, line_w=Pt(1.5))
txt(s, "Cold-Start (Simulated — restricted history)",
    Inches(6.9), Inches(4.38), Inches(5.9), Inches(0.4), size=12, bold=True, color=GREEN)
cold_rows = [
    ("History Size", "MF", "LightGCN", "Two-Tower"),
    ("3 interactions", "0.496", "0.400", "0.490 ✅"),
    ("Brand new (0)", "Cannot serve ✗", "Cannot serve ✗", "Works ✅"),
]
for ri, row in enumerate(cold_rows):
    for ci, cell in enumerate(row):
        lx = Inches(6.85) + ci*Inches(1.45)
        ly = Inches(4.8) + ri*Inches(0.65)
        fc = OVERLAY if ri==0 else (SURFACE if ri%2==0 else OVERLAY)
        box(s, lx, ly, Inches(1.4), Inches(0.58), fill=fc)
        color = GREEN if ("✅" in cell) else (RED if "✗" in cell else WHITE)
        txt(s, cell, lx+Inches(0.04), ly+Inches(0.06), Inches(1.35), Inches(0.48),
            size=10, bold=(ri==0), color=color, align=PP_ALIGN.CENTER)

# Insight
box(s, Inches(0.4), Inches(5.95), Inches(6.2), Inches(1.1), fill=OVERLAY, line=PURPLE, line_w=Pt(1.5))
txt(s, "Loss function > features: BPR on Two-Tower → HR@10 = 0.23 (collapsed!)\nInfoNCE is essential for dual-encoder training.",
    Inches(0.55), Inches(6.05), Inches(6.0), Inches(0.9), size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — ABLATION STUDY
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "12-Variant Ablation Study", "Isolate each component's contribution — change one thing at a time")

chart_data2 = ChartData()
chart_data2.categories = ["v1","v2","v3","v4","v4b","v4-BPR","v5","v5b","v5c","v6","v7","v8"]
chart_data2.add_series("HR@10", (0.6195,0.6210,0.6195,0.6355,0.6280,0.2295,
                                  0.6395,0.6385,0.6330,0.6355,0.6355,0.6305))

chart2 = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.3), Inches(1.3), Inches(7.5), Inches(4.8),
    chart_data2
).chart

chart2.has_title = False
chart2.font.color.rgb = WHITE
chart2.font.size = Pt(10)
try:
    chart2.chart_area.format.fill.solid()
    chart2.chart_area.format.fill.fore_color.rgb = SURFACE
    chart2.plot_area.format.fill.background()
except: pass

for series in chart2.plots[0].series:
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BLUE

# Annotations table
annotations = [
    ("v1", "Baseline InfoNCE b=256", "0.6195", SUBTEXT),
    ("v4", "+ Title text embeddings", "0.6355 ▲+2.6%", GREEN),
    ("v4-BPR", "BPR loss (wrong!)", "0.2295 ▼COLLAPSE", RED),
    ("v5", "+ GRU sequential", "0.6395 ▲ BEST", GREEN),
    ("v5b", "+ Rich text (noisy)", "0.6385 ▼", RED),
    ("v8", "FM-style gate weights", "0.6305", YELLOW),
]
box(s, Inches(7.9), Inches(1.3), Inches(5.1), Inches(5.8), fill=SURFACE)
txt(s, "Key Variants", Inches(8.0), Inches(1.38), Inches(4.9), Inches(0.4),
    size=13, bold=True, color=PURPLE)
for i, (ver, change, result, col) in enumerate(annotations):
    ly = Inches(1.88) + i*Inches(0.85)
    box(s, Inches(7.95), ly, Inches(5.0), Inches(0.75), fill=OVERLAY if i%2==0 else SURFACE)
    chip(s, ver, Inches(8.0), ly + Inches(0.1), w=Inches(0.8), h=Inches(0.28),
         fill=col, color=BG, size=9, bold=True)
    txt(s, change, Inches(8.9), ly+Inches(0.08), Inches(3.0), Inches(0.35), size=11, color=WHITE)
    txt(s, result, Inches(8.9), ly+Inches(0.43), Inches(3.0), Inches(0.28),
        size=10, bold=True, color=col)

txt(s, "Key finding: Loss function matters more than features.\nBPR → 0.23. InfoNCE → 0.64. Same model, different loss = 3×.",
    Inches(0.3), Inches(6.25), Inches(7.5), Inches(0.9), size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — INDUSTRY + SCALABILITY
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Industry Relevance & Scalability", "This is not an academic exercise — it's production architecture")

companies = [
    ("YouTube", "Two-Tower retrieval\n→ 1B+ videos, 2B users\nTop-1000 candidates\nthen deep ranker", PURPLE),
    ("Pinterest", "PinSage (graph-based)\nSame dual-encoder concept\nFAISS for pin retrieval\n2B+ pins served", BLUE),
    ("DoorDash", "Two-Tower for\nrestaurant + item recs\nReal-time user encoding\nFAISS candidate retrieval", GREEN),
    ("Airbnb", "Listing recommendations\nEmbedding similarity\nSame architecture\n$B in bookings driven", YELLOW),
    ("Spotify", "Podcast + song recs\nSequential encoding (like GRU)\nContrastive loss\nSame InfoNCE family", PEACH),
    ("Twitter / X", "Tweet ranking\nUser interest embeddings\nTwo-stage: retrieve → rank\nSame pattern", RED),
]
cw = Inches(2.1)
for i, (name, detail, col) in enumerate(companies):
    lx = Inches(0.3) + i*(cw + Inches(0.08))
    box(s, lx, Inches(1.3), cw, Inches(3.2), fill=SURFACE, line=col, line_w=Pt(2))
    txt(s, name, lx, Inches(1.35), cw, Inches(0.5),
        size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, detail, lx+Inches(0.1), Inches(1.9), cw-Inches(0.2), Inches(2.4),
        size=10, color=WHITE)

# Scalability
box(s, Inches(0.3), Inches(4.7), Inches(12.7), Inches(0.5), fill=PURPLE)
txt(s, "Why Two-Tower scales and LightGCN doesn't",
    Inches(0.5), Inches(4.75), Inches(12), Inches(0.38), size=14, bold=True, color=BG)

scale_points = [
    ("Two-Tower ✅", "Item vectors pre-computed once → FAISS. User tower is stateless.\nAdd more servers = linear scale. 29μs regardless of user count.\nWorks at YouTube's 2B users.", GREEN),
    ("LightGCN ✗", "Needs full adjacency matrix at inference.\n100M users × 100M items = doesn't fit in memory.\nCan't pre-compute static embeddings — graph changes.", RED),
    ("My System ✅", "Two-Tower handles retrieval + cold-start via FAISS.\nLightGCN could re-rank 1K candidates (small subgraph).\n34K queries/second on a single CPU core.", PURPLE),
]
for i, (title, body, col) in enumerate(scale_points):
    lx = Inches(0.3) + i*Inches(4.3)
    box(s, lx, Inches(5.35), Inches(4.1), Inches(1.85), fill=SURFACE, line=col, line_w=Pt(1.5))
    txt(s, title, lx+Inches(0.1), Inches(5.42), Inches(3.9), Inches(0.38),
        size=12, bold=True, color=col)
    txt(s, body, lx+Inches(0.1), Inches(5.82), Inches(3.9), Inches(1.2),
        size=10, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — KEY FINDINGS
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_header(s, "Key Findings", "What the data actually showed")

findings = [
    ("1", "Graph Structure Beats Features on Sparse Data",
     "LightGCN (0.729) beat all feature-enriched models.\nFG-LightGCN's gate confirms: 82% graph, 18% features.\nMulti-hop neighborhood captures 'users like you' — features can't.",
     BLUE, Inches(0.4), Inches(1.4)),
    ("2", "Text Helps, Rich Text Hurts",
     "Title embeddings: +2.6% (v4→v1).\nDescription + features text: −0.6% (v5b vs v5).\nNoise overwhelms signal on 99.97% sparse data.",
     GREEN, Inches(6.8), Inches(1.4)),
    ("3", "Loss Function > Architecture",
     "BPR on Two-Tower: HR@10 = 0.23 (collapsed).\nInfoNCE on same model: 0.64.\nSame architecture, wrong loss = 3× performance drop.",
     RED, Inches(0.4), Inches(4.3)),
    ("4", "No Single Best Model — Tradeoffs",
     "LightGCN: accuracy. Two-Tower: production + cold-start.\nFG-LightGCN: interpretable signal decomposition.\nThe right choice depends on the deployment constraint.",
     PURPLE, Inches(6.8), Inches(4.3)),
]
for num, title, body, col, lx, ty in findings:
    box(s, lx, ty, Inches(6.1), Inches(2.6), fill=SURFACE, line=col, line_w=Pt(2))
    chip(s, num, lx+Inches(0.1), ty+Inches(0.1), w=Inches(0.4), h=Inches(0.4),
         fill=col, color=BG, size=14, bold=True)
    txt(s, title, lx+Inches(0.65), ty+Inches(0.1), Inches(5.3), Inches(0.5),
        size=13, bold=True, color=col)
    txt(s, body, lx+Inches(0.15), ty+Inches(0.7), Inches(5.8), Inches(1.75),
        size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSE
# ══════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)

box(s, 0, 0, W, H, fill=BG)
box(s, 0, 0, W, Inches(0.08), fill=PURPLE)

txt(s, "Bottom Line", Inches(1.0), Inches(0.8), Inches(11), Inches(0.7),
    size=36, bold=True, color=PURPLE)

box(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.05), fill=PURPLE)

txt(s, "There is no single best model.\nLightGCN for accuracy.  Two-Tower for production.\nFeature-Gated LightGCN proves: on sparse data, the model itself learns to ignore features.",
    Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4), size=18, color=WHITE)

# 3 metric chips
for val, label, col, lx in [
    ("29 μs", "Per query (FAISS)", GREEN, Inches(1.2)),
    ("34K/sec", "Queries / second", PURPLE, Inches(4.8)),
    ("HR@10 0.729", "Best accuracy (LightGCN)", BLUE, Inches(8.4)),
]:
    box(s, lx, Inches(3.4), Inches(3.1), Inches(1.1), fill=SURFACE, line=col, line_w=Pt(2))
    txt(s, val, lx, Inches(3.45), Inches(3.1), Inches(0.55),
        size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, label, lx, Inches(3.98), Inches(3.1), Inches(0.4),
        size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)

txt(s, "github.com/nidhi1603/Two_Tower_Recommendation_System",
    Inches(1.0), Inches(4.75), Inches(11), Inches(0.4),
    size=13, color=BLUE, align=PP_ALIGN.CENTER)

url_shp = txt(s, "Live Demo →  https://two-tower-rec-sys.streamlit.app",
              Inches(1.0), Inches(5.2), Inches(11), Inches(0.4),
              size=14, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

txt(s, "Nidhi Rajani  ·  EAS 509  ·  Spring 2026",
    Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.4),
    size=12, color=SUBTEXT, align=PP_ALIGN.CENTER)

box(s, 0, H - Inches(0.08), W, Inches(0.08), fill=PURPLE)

# ── Save ──────────────────────────────────────────────────────
out = "/Users/nidhirajani/Desktop/Two_Tower_Recommendation_System/Recommendation_System_Presentation.pptx"
prs.save(out)
print(f"PPT saved → {out}")
